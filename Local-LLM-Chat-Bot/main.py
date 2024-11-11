
# Part 1
# from langchain_ollama import OllamaLLM
# from langchain_core.prompts import ChatPromptTemplate

# template = """
# Answer the question below.

# Here is the conversation history: {context}

# Question: {question}

# Answer:
# """

# model = OllamaLLM(model="llama3")
# prompt = ChatPromptTemplate.from_template(template)
# chain = prompt | model

# def handle_converstion():
#     context = ""
#     print("Welcome to the AI chatbot! Type 'exit' to quit.")
#     while True:
#         user_input = input("You: ")
#         if user_input.lower() == "exit":
#             break
#         result = chain.invoke({"context": context, "question": user_input})
#         print("Bot: ",result)
#         context += f"\nUser: {user_input}\nAI: {result}"
        

# if __name__ == "__main__":
#     handle_converstion()

# import os
# import fitz  # PyMuPDF
# import pandas as pd
# from pptx import Presentation
# import json
# from langchain_ollama import OllamaLLM
# from langchain_core.prompts import ChatPromptTemplate

# template = """
# Answer the question below.

# Here is the conversation history: {context}

# Question: {question}

# Answer:
# """

# model = OllamaLLM(model="llama3")
# prompt = ChatPromptTemplate.from_template(template)
# chain = prompt | model

# # Predefined folder path
# PREDEFINED_FOLDER_PATH = "/Users/sandrowtownsend/Documents/coding/Local LLM Chat Bot/data_base"

# def read_pdf(file_path):
#     doc = fitz.open(file_path)
#     text = ""
#     for page in doc:
#         text += page.get_text()
#     return text

# def read_txt(file_path):
#     with open(file_path, 'r') as file:
#         return file.read()

# def read_excel(file_path):
#     df = pd.read_excel(file_path)
#     return df.to_string()

# def read_ppt(file_path):
#     prs = Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text
#     return text

# def read_json(file_path):
#     with open(file_path, 'r') as file:
#         data = json.load(file)
#     return json.dumps(data, indent=4)

# def read_folder(folder_path):
#     context = ""
#     for filename in os.listdir(folder_path):
#         file_path = os.path.join(folder_path, filename)
#         if filename.endswith(".pdf"):
#             context += read_pdf(file_path)
#         elif filename.endswith(".txt"):
#             context += read_txt(file_path)
#         elif filename.endswith(".xlsx"):
#             context += read_excel(file_path)
#         elif filename.endswith(".pptx"):
#             context += read_ppt(file_path)
#         elif filename.endswith(".json"):
#             context += read_json(file_path)
#     return context

# def handle_conversation():
#     context = read_folder(PREDEFINED_FOLDER_PATH)
#     print("Welcome to the AI chatbot! Type 'exit' to quit.")
#     while True:
#         user_input = input("You: ")
#         if user_input.lower() == "exit":
#             break
#         result = chain.invoke({"context": context, "question": user_input})
#         print("Bot: ", result)
#         context += f"\nUser: {user_input}\nAI: {result}"

# if __name__ == "__main__":
#     handle_conversation()


    
import os
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
import json
from docx import Document
from langchain_ollama import OllamaLLM
from langchain_core.prompts import ChatPromptTemplate

template = """
Answer the question below.

Here is the conversation history: {context}

Question: {question}

Answer:
"""

model = OllamaLLM(model="llama3")
prompt = ChatPromptTemplate.from_template(template)
chain = prompt | model

# Predefined folder path
PREDEFINED_FOLDER_PATH = "/Users/sandrowtownsend/Documents/coding/Local-LLM-Chat-Bot/data_base"

def read_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def read_txt(file_path):
    with open(file_path, 'r') as file:
        return file.read()

def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()

def read_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def read_json(file_path):
    with open(file_path, 'r') as file:
        data = json.load(file)
    return json.dumps(data, indent=4)

def read_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def read_folder(folder_path):
    context = ""
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if filename.endswith(".pdf"):
            context += read_pdf(file_path)
        elif filename.endswith(".txt"):
            context += read_txt(file_path)
        elif filename.endswith(".xlsx"):
            context += read_excel(file_path)
        elif filename.endswith(".pptx"):
            context += read_ppt(file_path)
        elif filename.endswith(".json"):
            context += read_json(file_path)
        elif filename.endswith(".docx"):
            context += read_docx(file_path)
    return context

def handle_conversation():
    print("Welcome im your Local LLM Chat Bot, SandBot! Type 'exit' to quit.")
    while True:
        user_input = input("You: ")
        if user_input.lower() == "exit":
            break
        context = read_folder(PREDEFINED_FOLDER_PATH)
        result = chain.invoke({"context": context, "question": user_input})
        print("Bot: ", result)

if __name__ == "__main__":
    handle_conversation()

